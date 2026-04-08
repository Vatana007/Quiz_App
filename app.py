import streamlit as st
import json
import os
import uuid
import time
import PyPDF2
import docx
import pptx
from google import genai
from google.genai import types

# ==========================================
# 🔑 ដាក់ API KEY របស់អ្នកនៅទីនេះ
# ==========================================
MY_API_KEY = "ដាក់_API_KEY_របស់អ្នកនៅទីនេះ"

# --- System Setup ---
if not os.path.exists("quizzes"):
    os.makedirs("quizzes")

# --- MODERN ENTERPRISE UI TEMPLATE ---
st.set_page_config(page_title="ProQuiz AI", page_icon="⚡", layout="centered")

st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&family=Battambang:wght@400;700&display=swap');
    html, body, [class*="css"]  { font-family: 'Inter', 'Battambang', sans-serif; color: #1e293b; }
    .block-container { padding-top: 2rem; max-width: 750px; }
    
    .stButton>button { 
        width: 100%; border-radius: 6px; font-weight: 600; font-size: 15px;
        background-color: #0f172a; color: white; border: 1px solid #0f172a;
        padding: 0.6rem; transition: all 0.2s ease-in-out;
    }
    .stButton>button:hover { background-color: #334155; border-color: #334155; color: white; }
    
    .result-card { padding: 16px 20px; border-radius: 8px; margin-bottom: 16px; border: 1px solid #e2e8f0; background-color: #ffffff; box-shadow: 0 1px 3px rgba(0,0,0,0.04); }
    .correct-card { border-left: 4px solid #10b981; }
    .wrong-card { border-left: 4px solid #ef4444; }
    
    div.row-widget.stRadio > div{ background-color: #f8fafc; padding: 12px 16px; border-radius: 8px; border: 1px solid #cbd5e1; }
    
    .app-header { text-align: center; padding-bottom: 20px; border-bottom: 1px solid #e2e8f0; margin-bottom: 30px; }
    .app-title { font-weight: 700; font-size: 28px; color: #0f172a; margin-bottom: 5px; }
    .app-subtitle { color: #64748b; font-size: 15px; }

    /* --- លាក់ប៊ូតុង Streamlit Default Watermarks --- */
    [data-testid="stDecoration"] { display: none; }
    [data-testid="stToolbar"] { display: none; }
    [data-testid="manage-app-button"] { display: none !important; }
    div[class^="st-emotion-cache-"] > button { display: none !important; }
    </style>
""", unsafe_allow_html=True)

# --- CORE FUNCTIONS ---

# ប្រើប្រាស់ Cache សម្រាប់តែការអានឯកសារ ដើម្បីសន្សំសំចៃពេល
@st.cache_data(show_spinner=False)
def extract_text_from_upload(file_name, file_bytes, ext):
    text = ""
    try:
        if ext == 'txt': 
            text = file_bytes.decode("utf-8")
        elif ext == 'pdf':
            import io
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted: text += extracted + "\n"
        elif ext == 'docx':
            import io
            doc = docx.Document(io.BytesIO(file_bytes))
            for para in doc.paragraphs: text += para.text + "\n"
        elif ext == 'pptx':
            import io
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        return f"Error: {e}"
    return text.strip()

# មិនដាក់ Cache ទេ ដើម្បីឲ្យ AI បង្កើតសំណួរថ្មីប្លែកៗជានិច្ចពេលចុចម្តងៗ
def generate_quiz(text, num_questions, target_language):
    client = genai.Client(api_key=MY_API_KEY)
    
    prompt = f"""
    You are an expert AI teacher. Read the content and generate {num_questions} Multiple Choice Questions (MCQ).
    CRITICAL INSTRUCTION: Ensure the questions are highly diverse. Pick DIFFERENT details and concepts from the text than you might normally choose. Do not repeat standard questions.
    
    CRITICAL RULES:
    1. ALL generated questions, options, and correct answers MUST be written entirely in {target_language}.
    2. The value of "correct_answer" MUST exactly match one of the full strings inside the "options" array.
    
    Return STRICTLY as a JSON object matching this exact format:
    {{
      "questions": [
        {{
          "question": "Question text here?", 
          "options": ["A) Option 1", "B) Option 2", "C) Option 3", "D) Option 4"], 
          "correct_answer": "C) Option 3"
        }}
      ]
    }}
    
    Content to process:
    ---
    {text}
    ---
    """
    
    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = client.models.generate_content(
                model='gemini-2.5-flash',
                contents=prompt,
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    temperature=0.9 # បង្កើនភាពច្នៃប្រឌិតឲ្យចេញសំណួរថ្មី
                )
            )
            return json.loads(response.text)
        except Exception as e:
            error_msg = str(e)
            if "503" in error_msg or "UNAVAILABLE" in error_msg:
                if attempt < max_retries - 1:
                    time.sleep(4) # រង់ចាំ ៤ វិនាទី បើ Server កកស្ទះ
                    continue
            raise Exception(f"API Error: {error_msg}")

# --- ROUTING LOGIC ---
query_params = st.query_params
quiz_id = query_params.get("quiz_id")

# ==========================================
# ផ្ទាំងសិស្ស (STUDENT EXAM VIEW)
# ==========================================
if quiz_id:
    quiz_file_path = f"quizzes/{quiz_id}.json"
    
    if os.path.exists(quiz_file_path):
        with open(quiz_file_path, "r", encoding="utf-8") as f:
            saved_data = json.load(f)
            
        quiz_data = saved_data["data"]
        quiz_lang = saved_data["language"]
        
        # វចនានុក្រមប្រែភាសាស្វ័យប្រវត្តិសម្រាប់ UI
        ui = {
            "Khmer": {"title": "📝 ការប្រឡងសាកល្បង", "subtitle": "សូមជ្រើសរើសចម្លើយដែលត្រឹមត្រូវបំផុតសម្រាប់សំណួរនីមួយៗ", "submit": "បញ្ជូនចម្លើយ (Submit Assessment)", "result_title": "📊 លទ្ធផលរបស់អ្នក", "q": "សំណួរទី", "not_answered": "មិនបានឆ្លើយ", "correct_is": "ចម្លើយត្រូវគឺ", "correct": "ត្រឹមត្រូវ", "you_chose": "អ្នកជ្រើសរើស", "wrong": "ខុសហើយ", "your_choice": "ជម្រើសរបស់អ្នក", "passed": "ជាប់ (Passed)", "failed": "ធ្លាក់ (Failed)"},
            "English": {"title": "📝 Student Assessment", "subtitle": "Please select the best answer for each question.", "submit": "Submit Assessment", "result_title": "📊 Your Results", "q": "Question", "not_answered": "Not Answered", "correct_is": "Correct answer is", "correct": "Correct", "you_chose": "You chose", "wrong": "Incorrect", "your_choice": "Your choice", "passed": "Passed", "failed": "Failed"}
        }
        lang = ui[quiz_lang]

        st.markdown(f"<div class='app-header'><div class='app-title'>{lang['title']}</div><div class='app-subtitle'>{lang['subtitle']}</div></div>", unsafe_allow_html=True)
        
        with st.form(key="student_exam_form"):
            user_answers = {}
            for i, q in enumerate(quiz_data["questions"]):
                st.markdown(f"**{lang['q']} {i+1} :** {q['question']}")
                user_answers[i] = st.radio("Options:", q["options"], key=f"q_{i}", index=None, label_visibility="collapsed")
                st.write("") 
            
            submitted = st.form_submit_button(lang['submit'])
            
            if submitted:
                score = 0
                total_questions = len(quiz_data["questions"])
                st.markdown("---")
                st.markdown(f"<h3 style='text-align: center; color: #0f172a;'>{lang['result_title']}</h3>", unsafe_allow_html=True)
                
                for i, q in enumerate(quiz_data["questions"]):
                    correct = q["correct_answer"].strip()
                    user_ans = user_answers[i]
                    
                    if user_ans is None:
                        st.markdown(f'<div class="result-card wrong-card"><div style="font-weight: 600;">{lang["q"]} {i+1}: ❌ {lang["not_answered"]}</div><div style="color: #64748b; font-size: 14px; margin-top: 5px;">{lang["correct_is"]}: <b style="color:#0f172a;">{correct}</b></div></div>', unsafe_allow_html=True)
                        continue

                    user_ans = user_ans.strip()

                    if user_ans == correct:
                        score += 1
                        st.markdown(f'<div class="result-card correct-card"><div style="font-weight: 600; color: #10b981;">{lang["q"]} {i+1}: ✅ {lang["correct"]}</div><div style="color: #64748b; font-size: 14px; margin-top: 5px;">{lang["you_chose"]}: <span style="color:#0f172a;">{user_ans}</span></div></div>', unsafe_allow_html=True)
                    else:
                        st.markdown(f'<div class="result-card wrong-card"><div style="font-weight: 600; color: #ef4444;">{lang["q"]} {i+1}: ❌ {lang["wrong"]}</div><div style="color: #64748b; font-size: 14px; margin-top: 5px;">{lang["your_choice"]}: <del>{user_ans}</del></div><div style="color: #64748b; font-size: 14px;">{lang["correct_is"]}: <b style="color:#0f172a;">{correct}</b></div></div>', unsafe_allow_html=True)
                
                percentage = int((score / total_questions) * 100)
                st.markdown("<br>", unsafe_allow_html=True)
                
                if percentage >= 50:
                    st.success(f"🎉 **{lang['passed']}!** Score: {score} / {total_questions} ({percentage}%)")
                else:
                    st.error(f"📉 **{lang['failed']}!** Score: {score} / {total_questions} ({percentage}%)")
    else:
        st.warning("⚠️ Quiz not found. The link may be invalid or expired.")

# ==========================================
# ផ្ទាំងគ្រូ (TEACHER / ADMIN VIEW)
# ==========================================
else:
    st.markdown("""
        <div class="app-header">
            <div class="app-title">⚡ ProQuiz Generator</div>
            <div class="app-subtitle">Automated assessment generation for educators</div>
        </div>
    """, unsafe_allow_html=True)

    with st.sidebar:
        st.markdown("### ⚙️ Configuration")
        selected_language = st.selectbox("🗣️ Output Language", ("Khmer", "English"))
        num_q = st.number_input("📝 Number of Questions", min_value=1, max_value=50, value=5)
        st.markdown("---")

    st.markdown("#### 📄 Upload Source Material")
    uploaded_file = st.file_uploader("", type=['pdf', 'docx', 'pptx', 'txt'], label_visibility="collapsed")

    if uploaded_file:
        st.info(f"File attached: **{uploaded_file.name}**")
        
        if st.button("Generate Assessment Link"):
            if MY_API_KEY == "ដាក់_API_KEY_របស់អ្នកនៅទីនេះ":
                st.error("⚠️ សូមដូរពាក្យ 'ដាក់_API_KEY_របស់អ្នកនៅទីនេះ' នៅបន្ទាត់ទី ១៥ ទៅជា API Key របស់អ្នកពិតប្រាកដសិន។")
            else:
                with st.spinner("Analyzing document and generating questions..."):
                    # Extract Data (Cached)
                    file_bytes = uploaded_file.getvalue()
                    ext = uploaded_file.name.split('.')[-1].lower()
                    document_text = extract_text_from_upload(uploaded_file.name, file_bytes, ext)
                    
                    if document_text and not document_text.startswith("Error:"):
                        try:
                            # Generate Quiz (Not Cached - New questions every time)
                            quiz_data = generate_quiz(document_text, num_q, selected_language)
                            unique_id = str(uuid.uuid4())
                            
                            save_payload = {"language": selected_language, "data": quiz_data}
                            with open(f"quizzes/{unique_id}.json", "w", encoding="utf-8") as f:
                                json.dump(save_payload, f, ensure_ascii=False, indent=4)
                            
                            # យក URL ជាក់ស្តែងរបស់កុំព្យូទ័រ/Server ដើម្បីងាយស្រួល Copy
                            base_url = "http://localhost:8501" 
                            shareable_link = f"{base_url}/?quiz_id={unique_id}"
                            
                            st.success("✅ Assessment generated successfully!")
                            st.markdown(f"""
                            <div style="background-color: #f8fafc; padding: 20px; border-radius: 8px; border: 1px dashed #cbd5e1; text-align: center;">
                                <p style="color: #475569; font-weight: 500; font-size: 14px;">Copy and share this link with your students:</p>
                                <code style="font-size: 15px; color: #2563eb; background: #fff; padding: 8px 12px; border: 1px solid #e2e8f0; border-radius: 4px; display: inline-block; margin: 10px 0;">{shareable_link}</code>
                                <br><br>
                                <a href="{shareable_link}" target="_blank" style="text-decoration: none; background-color: #0f172a; color: white; padding: 8px 16px; border-radius: 6px; font-weight: 600; font-size: 13px;">Open Student View</a>
                            </div>
                            """, unsafe_allow_html=True)

                        except Exception as e:
                            st.error(f"❌ Generation failed: {e}")
                    else:
                        st.error(f"❌ Could not extract readable text: {document_text}")